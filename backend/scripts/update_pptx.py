from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def replace_text_in_presentation(pptx_path, output_path, replacements):
    prs = Presentation(pptx_path)
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num in replacements:
            rep = replacements[slide_num]
            
            # Handle Title if provided
            if "title" in rep and slide.has_notes_slide == False or True: # Force check all
                if slide.shapes.title:
                    slide.shapes.title.text = rep["title"]
            
            # Handle specific text replacements within shapes
            if "text_reps" in rep:
                for target, replacement in rep["text_reps"].items():
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if target in run.text:
                                        run.text = run.text.replace(target, replacement)
            
            # Handle list replacement (if we want to just overwrite all text frames that are not titles)
            if "body_list" in rep:
                body_idx = 0
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        # Only replace non-title text frames that have content
                        if shape.text_frame.text.strip() and body_idx < len(rep["body_list"]):
                            shape.text_frame.text = rep["body_list"][body_idx]
                            body_idx += 1

    prs.save(output_path)

# Define replacements based on research paper mapping
replacements = {
    1: {
        "title": "Intelligent Regulatory Insight System (I.R.I.S.)",
        "text_reps": {"Project I.R.I.S.": "AI Driven Financial Forensics and Compliance Platform"}
    },
    9: {
        "title": "Our Mission & Vision",
        "body_list": [
            "Mission: To support fraud detection, risk assessment, and regulatory compliance for Indian public firms.",
            "Vision: A unified AI architecture that operates with minimal human intervention and near real-time responsiveness.",
            "Objectives: Automate ingestion, forensic metric computation, and composite risk scoring."
        ]
    },
    11: {
        "title": "Financial Risk Dimensions",
        "body_list": [
            "Earnings Quality", "Leverage & Solvency", "Liquidity Pressure", "Revenue Stability", 
            "Governance & Disclosure", "Market Behaviour Anomalies"
        ]
    },
    12: {
        "title": "PESTLE Analysis",
        "text_reps": {
            "Political": "P: SEBI Strategic Shift to Tech-driven Supervision",
            "Economical": "E: Impact on Capital Markets & Economic Stability",
            "Social": "S: Improving Public Confidence in Financial Integrity",
            "Technological": "T: Multi-Agent AI Systems & Language Models",
            "Legal": "L: Compliance with Digital Forensics & SEBI Circulars",
            "Environmental": "E: Sustainable & Fair Market Ecosystems"
        }
    },
    16: {
        "title": "Multi-Agent Architecture",
        "body_list": [
            "Agent 1: Multi-source Data Ingestion (Yahoo/NSE/BSE)",
            "Agent 2: Computation of 29 Forensic Metrics",
            "Agent 3: Automated Composite Risk Scoring",
            "Agent 4: Regulatory Compliance Validation",
            "Agent 5: Human-Readable Report Generation"
        ]
    },
    23: {
        "title": "Integrated Forensic Services",
        "body_list": [
            "1. Automated Forensic Metric Computation (Altman, Beneish, Benford)",
            "2. FinBERT Sentiment Analysis on News & Filings",
            "3. Peer Benchmarking & Z-Score Analysis",
            "4. Regulatory Event Tracking (SEBI Orders/Circulars)",
            "5. Graph-based RPT Visualization"
        ]
    },
    33: {
        "title": "User Segments & Personas",
        "text_reps": {
            "Anna Wilson": "SEBI Analysts & Regulators",
            "Lawyer": "Forensic Auditors",
            "Hobby 1": "Institutional Investors",
            "Hobby 2": "Compliance Officers"
        },
        "body_list": [
            "Needs: Traceable data lineage, interpretable metrics, and explainable risk scores.",
            "Goals: Rapid fraud detection and evidence-based enforcement."
        ]
    },
    45: {
        "title": "IRS Maturity Model",
        "body_list": [
            "Level 1: Standard Rule-based Compliance Checks",
            "Level 2: Advanced AI Sentiment & Quantitative Scores",
            "Level 3: Autonomous Multi-Agent Orchestration & Q&A"
        ]
    },
    48: {
        "title": "The Research Team",
        "text_reps": {
            "Jenna Doe": "Swapna Lokande & Rishikesh Koli",
            "Timmy Jimmy": "Tejaswini Patil, Isha Sonkusare & Yashraj Kulkrani"
        }
    }
}

if __name__ == "__main__":
    input_pptx = "Copy of Financial Markets Consulting Services by Slidesgo.pptx"
    output_pptx = "Financial_Markets_Consulting_IRS_Final.pptx"
    
    replace_text_in_presentation(input_pptx, output_pptx, replacements)
    print(f"Successfully updated presentation saved to {output_pptx}")
