import sys
from pptx import Presentation

def extract_pptx_content(pptx_path):
    prs = Presentation(pptx_path)
    content = []
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "title": "",
            "shapes": [],
            "notes": ""
        }
        
        if slide.has_notes_slide:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text
            
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if shape == slide.shapes.title:
                        slide_info["title"] = text
                    else:
                        slide_info["shapes"].append(text)
        
        content.append(slide_info)
    
    return content

if __name__ == "__main__":
    path = "Copy of Financial Markets Consulting Services by Slidesgo.pptx"
    if len(sys.argv) > 1:
        path = sys.argv[1]
    
    try:
        slides = extract_pptx_content(path)
        for s in slides:
            print(f"--- Slide {s['slide_number']}: {s['title']} ---")
            for shape_text in s['shapes']:
                print(f"  - {shape_text}")
            if s['notes']:
                print(f"  Notes: {s['notes']}")
            print("\n")
    except Exception as e:
        print(f"Error: {e}")
